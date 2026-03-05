import base64
import json
import math
import os
import sys
import tempfile
from typing import Any, Dict, Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.parts.image import Image as PptxImage
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400

DEFAULT_THEME: Dict[str, Any] = {
    "name": "Emerald Pro",
    "colors": {
        "background": "#FFFFFF",
        "primary": "#0F9D58",
        "secondary": "#1B1F24",
        "accent": "#00C2A8",
        "muted": "#6B7280",
        "light": "#F3F4F6",
        "dark": "#111827",
        "card": "#F8FAFC",
    },
    "fonts": {
        "title": "Aptos Display",
        "body": "Aptos",
        "mono": "Consolas",
    },
    "sizes": {
        "title": 40,
        "subtitle": 24,
        "section": 44,
        "body": 18,
        "small": 14,
        "stat": 32,
        "quote": 36,
    },
    "line_spacing": 1.15,
    "slide": {
        "width": 13.333,
        "height": 7.5,
    },
    "margins": {
        "left": 0.6,
        "right": 0.6,
        "top": 0.5,
        "bottom": 0.5,
    },
    "grid": {
        "gutter": 0.3,
    },
    "footer": {
        "show": True,
        "text": "",
        "size": 12,
    },
}


def deep_merge(base: Dict[str, Any], override: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not override:
        return base
    result = dict(base)
    for key, value in override.items():
        if isinstance(value, dict) and isinstance(result.get(key), dict):
            result[key] = deep_merge(result[key], value)
        else:
            result[key] = value
    return result


def normalize_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (str, int, float, bool)):
        return str(value)
    if isinstance(value, dict):
        for key in ("text", "title", "label", "value", "name", "content"):
            if key in value and value[key]:
                return str(value[key])
        return json.dumps(value, ensure_ascii=False)
    if isinstance(value, list):
        return ", ".join([normalize_text(item) for item in value if item is not None])
    return str(value)


def normalize_bullets(items: Iterable[Any]) -> List[str]:
    bullets: List[str] = []
    for item in items or []:
        if isinstance(item, list):
            for nested in item:
                text = normalize_text(nested)
                if text:
                    bullets.append(text)
        else:
            text = normalize_text(item)
            if text:
                bullets.append(text)
    return bullets


def ensure_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()


def to_rgb(color: str, fallback: str = "#000000") -> RGBColor:
    value = (color or fallback).strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join([c * 2 for c in value])
    if len(value) != 6:
        value = fallback.replace("#", "")
    r = int(value[0:2], 16)
    g = int(value[2:4], 16)
    b = int(value[4:6], 16)
    return RGBColor(r, g, b)


def inches(value: float) -> Inches:
    return Inches(float(value))


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_name: str,
    font_size: int,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: Optional[float] = None,
):
    box = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = normalize_text(text)
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = ensure_run(p)
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = to_rgb(color)
    return box


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: List[Any],
    font_name: str,
    font_size: int,
    color: str,
    line_spacing: float,
    level: int = 0,
):
    box = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_items = normalize_bullets(bullets)
    for idx, bullet in enumerate(bullet_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = ensure_run(p)
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = to_rgb(color)
    return box


def add_accent_bar(slide, left: float, top: float, width: float, height: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, inches(left), inches(top), inches(width), inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(color)
    shape.line.fill.background()
    return shape


def set_background(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = to_rgb(color)


def image_size_in_inches(path: str) -> Tuple[float, float]:
    image = PptxImage.from_file(path)
    width, height = image.extents
    return width / EMU_PER_INCH, height / EMU_PER_INCH


def save_temp_image(data_base64: str) -> str:
    decoded = base64.b64decode(data_base64)
    fd, path = tempfile.mkstemp(suffix=".png")
    with os.fdopen(fd, "wb") as f:
        f.write(decoded)
    return path


def add_picture_fit(
    slide,
    path: str,
    left: float,
    top: float,
    width: float,
    height: float,
    mode: str = "contain",
):
    img_w, img_h = image_size_in_inches(path)
    box_ratio = width / height
    img_ratio = img_w / img_h if img_h else box_ratio
    if mode == "cover":
        pic = slide.shapes.add_picture(path, inches(left), inches(top), inches(width), inches(height))
        if img_ratio > box_ratio:
            crop = (img_ratio - box_ratio) / (2 * img_ratio)
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ratio < box_ratio:
            crop = (box_ratio - img_ratio) / (2 * box_ratio)
            pic.crop_top = crop
            pic.crop_bottom = crop
        return pic
    scale = min(width / img_w, height / img_h) if img_w and img_h else 1.0
    new_w = img_w * scale
    new_h = img_h * scale
    x = left + (width - new_w) / 2
    y = top + (height - new_h) / 2
    return slide.shapes.add_picture(path, inches(x), inches(y), inches(new_w), inches(new_h))


def add_footer(slide, theme: Dict[str, Any], text: str):
    if not theme.get("footer", {}).get("show", True):
        return
    footer_text = text or theme.get("footer", {}).get("text", "")
    if not footer_text:
        return
    width = theme["slide"]["width"]
    left = theme["margins"]["left"]
    right = theme["margins"]["right"]
    height = 0.3
    top = theme["slide"]["height"] - theme["margins"]["bottom"] + 0.05
    add_textbox(
        slide,
        left,
        top,
        width - left - right,
        height,
        footer_text,
        theme["fonts"]["body"],
        theme["footer"]["size"],
        theme["colors"]["muted"],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_notes(slide, notes: Optional[str]):
    if not notes:
        return
    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_title_block(
    slide,
    theme: Dict[str, Any],
    title: str,
    subtitle: Optional[str],
    section: Optional[str] = None,
):
    left = theme["margins"]["left"]
    top = theme["margins"]["top"]
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    add_textbox(
        slide,
        left,
        top,
        width,
        0.9,
        title or "",
        theme["fonts"]["title"],
        theme["sizes"]["title"],
        theme["colors"]["secondary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )
    if subtitle:
        add_textbox(
            slide,
            left,
            top + 0.85,
            width,
            0.6,
            subtitle,
            theme["fonts"]["body"],
            theme["sizes"]["subtitle"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
    if section:
        add_textbox(
            slide,
            left,
            top - 0.25,
            width,
            0.3,
            section,
            theme["fonts"]["body"],
            theme["sizes"]["small"],
            theme["colors"]["accent"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
        add_accent_bar(slide, left, top + 1.55, 1.0, 0.05, theme["colors"]["accent"])


def add_cover_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    left = theme["margins"]["left"]
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    top = 2.0
    add_textbox(
        slide,
        left,
        top,
        width,
        1.2,
        slide_data.get("title", ""),
        theme["fonts"]["title"],
        theme["sizes"]["section"],
        theme["colors"]["secondary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    subtitle = slide_data.get("subtitle")
    if subtitle:
        add_textbox(
            slide,
            left,
            top + 1.0,
            width,
            0.7,
            subtitle,
            theme["fonts"]["body"],
            theme["sizes"]["subtitle"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
    meta = slide_data.get("meta")
    if meta:
        add_textbox(
            slide,
            left,
            theme["slide"]["height"] - 1.2,
            width,
            0.4,
            meta,
            theme["fonts"]["body"],
            theme["sizes"]["small"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_section_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_textbox(
        slide,
        theme["margins"]["left"],
        2.4,
        theme["slide"]["width"] - theme["margins"]["left"] - theme["margins"]["right"],
        1.0,
        slide_data.get("title", ""),
        theme["fonts"]["title"],
        theme["sizes"]["section"],
        theme["colors"]["primary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    subtitle = slide_data.get("subtitle")
    if subtitle:
        add_textbox(
            slide,
            theme["margins"]["left"],
            3.4,
            theme["slide"]["width"] - theme["margins"]["left"] - theme["margins"]["right"],
            0.6,
            subtitle,
            theme["fonts"]["body"],
            theme["sizes"]["subtitle"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
        )
    add_accent_bar(slide, theme["margins"]["left"], 4.2, 2.5, 0.08, theme["colors"]["accent"])
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_content_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    bullets = slide_data.get("content") or slide_data.get("bullets") or []
    columns = slide_data.get("columns", 1)
    left = theme["margins"]["left"]
    top = 2.1
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    if columns == 2 and len(bullets) > 1:
        mid = left + width / 2 + theme["grid"]["gutter"] / 2
        left_col = bullets[: math.ceil(len(bullets) / 2)]
        right_col = bullets[math.ceil(len(bullets) / 2) :]
        add_bullets(
            slide,
            left,
            top,
            width / 2 - theme["grid"]["gutter"],
            height,
            left_col,
            theme["fonts"]["body"],
            theme["sizes"]["body"],
            theme["colors"]["secondary"],
            theme["line_spacing"],
        )
        add_bullets(
            slide,
            mid,
            top,
            width / 2 - theme["grid"]["gutter"],
            height,
            right_col,
            theme["fonts"]["body"],
            theme["sizes"]["body"],
            theme["colors"]["secondary"],
            theme["line_spacing"],
        )
    else:
        add_bullets(
            slide,
            left,
            top,
            width,
            height,
            bullets,
            theme["fonts"]["body"],
            theme["sizes"]["body"],
            theme["colors"]["secondary"],
            theme["line_spacing"],
        )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_two_column_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    left = theme["margins"]["left"]
    top = 2.1
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    gap = theme["grid"]["gutter"]
    column_width = (width - gap) / 2
    add_bullets(
        slide,
        left,
        top,
        column_width,
        height,
        slide_data.get("left", []),
        theme["fonts"]["body"],
        theme["sizes"]["body"],
        theme["colors"]["secondary"],
        theme["line_spacing"],
    )
    add_bullets(
        slide,
        left + column_width + gap,
        top,
        column_width,
        height,
        slide_data.get("right", []),
        theme["fonts"]["body"],
        theme["sizes"]["body"],
        theme["colors"]["secondary"],
        theme["line_spacing"],
    )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_image_text_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any], image_left: bool):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    left = theme["margins"]["left"]
    top = 2.1
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    gap = theme["grid"]["gutter"]
    image_box = (width - gap) * 0.45
    text_box = width - gap - image_box
    image_x = left if image_left else left + text_box + gap
    text_x = left + image_box + gap if image_left else left
    image = slide_data.get("image", {})
    image_path = image.get("path")
    image_base64 = image.get("base64")
    if image_base64 and not image_path:
        image_path = save_temp_image(image_base64)
        slide_data.setdefault("_temp_files", []).append(image_path)
    if image_path and os.path.exists(image_path):
        add_picture_fit(
            slide,
            image_path,
            image_x,
            top,
            image_box,
            height,
            image.get("fit", "cover"),
        )
    bullets = slide_data.get("content") or slide_data.get("bullets") or []
    add_bullets(
        slide,
        text_x,
        top,
        text_box,
        height,
        bullets,
        theme["fonts"]["body"],
        theme["sizes"]["body"],
        theme["colors"]["secondary"],
        theme["line_spacing"],
    )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_quote_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    quote = slide_data.get("quote", "")
    author = slide_data.get("author", "")
    left = theme["margins"]["left"]
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    add_textbox(
        slide,
        left,
        2.0,
        width,
        2.0,
        quote,
        theme["fonts"]["title"],
        theme["sizes"]["quote"],
        theme["colors"]["secondary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if author:
        add_textbox(
            slide,
            left,
            4.0,
            width,
            0.5,
            author,
            theme["fonts"]["body"],
            theme["sizes"]["subtitle"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
    add_accent_bar(slide, left, 1.7, 1.2, 0.06, theme["colors"]["accent"])
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_stats_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    items = slide_data.get("items", [])
    left = theme["margins"]["left"]
    top = 2.2
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    count = max(1, min(4, len(items)))
    card_width = (width - theme["grid"]["gutter"] * (count - 1)) / count
    for idx in range(count):
        item = items[idx]
        card_left = left + idx * (card_width + theme["grid"]["gutter"])
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            inches(card_left),
            inches(top),
            inches(card_width),
            inches(height * 0.7),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = to_rgb(theme["colors"]["card"])
        card.line.fill.background()
        add_textbox(
            slide,
            card_left + 0.2,
            top + 0.3,
            card_width - 0.4,
            0.6,
            item.get("label", ""),
            theme["fonts"]["body"],
            theme["sizes"]["small"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
        add_textbox(
            slide,
            card_left + 0.2,
            top + 0.8,
            card_width - 0.4,
            0.8,
            item.get("value", ""),
            theme["fonts"]["title"],
            theme["sizes"]["stat"],
            theme["colors"]["secondary"],
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
        delta = item.get("delta")
        if delta:
            add_textbox(
                slide,
                card_left + 0.2,
                top + 1.6,
                card_width - 0.4,
                0.4,
                delta,
                theme["fonts"]["body"],
                theme["sizes"]["small"],
                theme["colors"]["accent"],
                align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP,
            )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_timeline_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    items = slide_data.get("items", [])
    left = theme["margins"]["left"]
    top = 2.4
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    row_height = 0.9
    for idx, item in enumerate(items[:5]):
        y = top + idx * row_height
        add_accent_bar(slide, left, y + 0.2, 0.15, 0.15, theme["colors"]["accent"])
        add_textbox(
            slide,
            left + 0.3,
            y,
            2.0,
            0.5,
            item.get("title", ""),
            theme["fonts"]["body"],
            theme["sizes"]["body"],
            theme["colors"]["secondary"],
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
        add_textbox(
            slide,
            left + 2.4,
            y,
            width - 2.4,
            0.6,
            item.get("text", ""),
            theme["fonts"]["body"],
            theme["sizes"]["body"],
            theme["colors"]["muted"],
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_table_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    columns = slide_data.get("columns", [])
    rows = slide_data.get("rows", [])
    if not columns:
        fallback = dict(slide_data)
        fallback["content"] = [", ".join([normalize_text(value) for value in row]) for row in rows] or [
            "No data provided"
        ]
        add_content_slide(prs, fallback, theme)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    left = theme["margins"]["left"]
    top = 2.2
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    col_count = len(columns)
    row_count = len(rows) + 1
    table_shape = slide.shapes.add_table(row_count, col_count, inches(left), inches(top), inches(width), inches(height))
    table = table_shape.table
    for col_idx, col_name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = to_rgb(theme["colors"]["light"])
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.name = theme["fonts"]["body"]
                run.font.size = Pt(theme["sizes"]["small"])
                run.font.bold = True
                run.font.color.rgb = to_rgb(theme["colors"]["secondary"])
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                if paragraph.runs:
                    run = paragraph.runs[0]
                    run.font.name = theme["fonts"]["body"]
                    run.font.size = Pt(theme["sizes"]["small"])
                    run.font.color.rgb = to_rgb(theme["colors"]["secondary"])
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def add_chart_slide(prs, slide_data: Dict[str, Any], theme: Dict[str, Any]):
    chart_spec = slide_data.get("chart", {})
    if not chart_spec.get("series"):
        fallback = dict(slide_data)
        fallback["content"] = slide_data.get("content") or ["No chart data provided"]
        add_content_slide(prs, fallback, theme)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, slide_data.get("background", theme["colors"]["background"]))
    add_title_block(
        slide,
        theme,
        slide_data.get("title", ""),
        slide_data.get("subtitle"),
        slide_data.get("section"),
    )
    chart_type = chart_spec.get("type", "bar").lower()
    chart_data = ChartData()
    categories = chart_spec.get("categories", [])
    chart_data.categories = categories
    for series in chart_spec.get("series", []):
        chart_data.add_series(series.get("name", ""), series.get("values", []))
    left = theme["margins"]["left"]
    top = 2.1
    width = theme["slide"]["width"] - left - theme["margins"]["right"]
    height = theme["slide"]["height"] - top - theme["margins"]["bottom"] - 0.2
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart_shape = slide.shapes.add_chart(
        chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
        inches(left),
        inches(top),
        inches(width),
        inches(height),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    add_notes(slide, slide_data.get("notes"))
    add_footer(slide, theme, slide_data.get("footer", ""))


def expand_slides(slides: List[Dict[str, Any]], max_bullets: int) -> List[Dict[str, Any]]:
    expanded = []
    for slide in slides:
        if slide.get("type") in (None, "content") and slide.get("allow_split", True):
            bullets = slide.get("content") or slide.get("bullets") or []
            if len(bullets) > max_bullets:
                chunks = [bullets[i : i + max_bullets] for i in range(0, len(bullets), max_bullets)]
                for index, chunk in enumerate(chunks):
                    new_slide = dict(slide)
                    new_slide["content"] = chunk
                    if index > 0 and new_slide.get("title"):
                        new_slide["title"] = f"{new_slide['title']} ({index + 1})"
                    expanded.append(new_slide)
                continue
        expanded.append(slide)
    return expanded


def build_presentation(data: Dict[str, Any]) -> Tuple[Presentation, List[str]]:
    theme = deep_merge(DEFAULT_THEME, data.get("theme"))
    prs = Presentation()
    prs.slide_width = inches(theme["slide"]["width"])
    prs.slide_height = inches(theme["slide"]["height"])
    temp_files: List[str] = []
    slides = data.get("slides", [])
    slides = expand_slides(slides, data.get("max_bullets_per_slide", 8))
    for slide in slides:
        slide.setdefault("_temp_files", [])
        slide_type = (slide.get("type") or "content").lower()
        if slide_type == "title":
            add_cover_slide(prs, slide, theme)
        elif slide_type == "section":
            add_section_slide(prs, slide, theme)
        elif slide_type == "two_column":
            add_two_column_slide(prs, slide, theme)
        elif slide_type == "image_left":
            add_image_text_slide(prs, slide, theme, True)
        elif slide_type == "image_right":
            add_image_text_slide(prs, slide, theme, False)
        elif slide_type == "quote":
            add_quote_slide(prs, slide, theme)
        elif slide_type == "stats":
            add_stats_slide(prs, slide, theme)
        elif slide_type == "timeline":
            add_timeline_slide(prs, slide, theme)
        elif slide_type == "table":
            add_table_slide(prs, slide, theme)
        elif slide_type == "chart":
            add_chart_slide(prs, slide, theme)
        else:
            add_content_slide(prs, slide, theme)
        temp_files.extend(slide.get("_temp_files", []))
    return prs, temp_files


def normalize_legacy_schema(data: Dict[str, Any]) -> Dict[str, Any]:
    slides = data.get("slides") or []
    normalized = []
    for slide in slides:
        if "content" in slide or "title" in slide:
            normalized.append(slide)
        else:
            normalized.append({"type": "content", "title": slide.get("heading", ""), "content": slide.get("items", [])})
    if data.get("title") and not slides:
        normalized.append({"type": "title", "title": data.get("title"), "subtitle": data.get("subtitle")})
    data["slides"] = normalized
    return data


def main():
    raw = sys.stdin.read()
    payload = json.loads(raw) if raw else {}
    data = normalize_legacy_schema(payload)
    output_file = data.get("output_file") or data.get("output")
    if not output_file:
        output_file = os.path.abspath("presentation.pptx")
    prs, temp_files = build_presentation(data)
    prs.save(output_file)
    for path in temp_files:
        try:
            os.remove(path)
        except OSError:
            pass
    print(output_file)


if __name__ == "__main__":
    main()
